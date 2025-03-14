import google.generativeai as genai
import streamlit as st
from pptx import Presentation
import tempfile
import os
import pandas as pd
import time
from google.api_core import retry
from google.api_core.exceptions import ResourceExhausted  # Import correct exception

# Configure Gemini API
genai.configure(api_key=st.secrets["GEMINI_API_KEY"])

# Custom retry configuration
def if_retryable(error):
    return isinstance(error, (genai.types.BlockedPromptException, 
                                genai.types.StopCandidateException,
                                ResourceExhausted))  # Use imported ResourceExhausted

retry_policy = retry.Retry(
    predicate=if_retryable,
    initial=1,
    maximum=10,
    multiplier=2,
    deadline=300
)

# Session state management
if 'processed' not in st.session_state:
    st.session_state.processed = {
        'file_hash': None,
        'slides': [],
        'flashcards': [],
        'api_calls': 0,
        'last_call': 0
    }

# Rate limiting configuration
MAX_CALLS_PER_MINUTE = 30  # Adjust based on your quota
MIN_INTERVAL = max(60 / MAX_CALLS_PER_MINUTE, 1)

@retry_policy
def safe_generate(prompt):
    current_time = time.time()
    elapsed = current_time - st.session_state.processed['last_call']
    
    if elapsed < MIN_INTERVAL:
        time.sleep(MIN_INTERVAL - elapsed)
    
    model = genai.GenerativeModel("gemini-1.5-pro-latest")
    response = model.generate_content(prompt)
    
    st.session_state.processed['api_calls'] += 1
    st.session_state.processed['last_call'] = time.time()
    
    return response.text

def process_ppt(file):
    """Process PowerPoint file with caching"""
    file_hash = hash(file.getvalue())
    if file_hash == st.session_state.processed['file_hash']:
        return st.session_state.processed['slides']
    
    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        tmp.write(file.read())
        prs = Presentation(tmp.name)
    
    results = []
    for i, slide in enumerate(prs.slides):
        text = "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)
        results.append({
            "number": i+1,
            "content": text,
            "analysis": "",
            "flashcards": []
        })
    
    os.unlink(tmp.name)
    st.session_state.processed['file_hash'] = file_hash
    st.session_state.processed['slides'] = results
    return results

def generate_analysis(content):
    """Generate slide analysis with efficient prompt"""
    prompt = f"Provide concise exam-focused analysis in bullet points:\n{content}"
    return safe_generate(prompt)

def generate_flashcards(content, card_type):
    """Generate formatted flashcards"""
    formats = {
        "Q&A": "Generate 2-3 Q&A flashcards using:\nQ: question\nA: answer",
        "Term/Definition": "Create 2-3 term/definition pairs using:\nT: term\nD: definition",
        "MCQs": "Generate 2 MCQs with 4 options using:\nQ: question\nOptions: a) b) c) d)\nA: answer"
    }
    prompt = f"{formats[card_type]} from:\n{content}"
    return safe_generate(prompt)

def parse_flashcards(text, card_type):
    """Parse generated flashcards into structured format"""
    cards = []
    current = {}
    lines = text.split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if card_type == "MCQs":
            if line.startswith('Q:'):
                current = {'question': line[3:].strip()}
            elif line.startswith('Options:'):
                current['options'] = line[8:].strip()
            elif line.startswith('A:'):
                current['answer'] = line[3:].strip()
                cards.append(current)
        else:
            key = 'term' if card_type == "Term/Definition" else 'question'
            val = 'definition' if card_type == "Term/Definition" else 'answer'
            
            if line.startswith(f"{key[0].upper()}:"):
                current = {key: line[3:].strip()}
            elif line.startswith(f"{val[0].upper()}:"):
                current[val] = line[3:].strip()
                cards.append(current)
    
    return cards

# Streamlit UI
st.set_page_config(page_title="AI PPT Summarizer", layout="wide")
st.title("📚 AI PPT Summarizer with Smart Features")

with st.sidebar:
    st.header("Settings ⚙️")
    card_type = st.radio("Flashcard Style", ["Q&A", "Term/Definition", "MCQs"])
    analysis_depth = st.select_slider("Analysis Depth", ["Concise", "Normal", "Detailed"])
    st.progress(st.session_state.processed['api_calls'] % MAX_CALLS_PER_MINUTE / MAX_CALLS_PER_MINUTE)
    st.caption(f"API Calls: {st.session_state.processed['api_calls']} (this session)")

uploaded_file = st.file_uploader("Upload PowerPoint file", type=["ppt", "pptx"])

if uploaded_file:
    slides = process_ppt(uploaded_file)
    progress_bar = st.progress(0)
    
    for idx, slide in enumerate(slides):
        col1, col2 = st.columns([1, 2])
        
        with col1:
            with st.expander(f"Slide {slide['number']} Content", expanded=True):
                st.write(slide['content'] or "No text detected")
        
        with col2:
            tab1, tab2 = st.tabs(["Analysis 📝", "Flashcards 🗂️"])
            
            with tab1:
                if not slide['analysis']:
                    try:
                        slide['analysis'] = generate_analysis(slide['content'])
                    except Exception as e:
                        st.error(f"Analysis failed: {str(e)}")
                st.markdown(slide['analysis'])
            
            with tab2:
                if not slide['flashcards']:
                    try:
                        raw = generate_flashcards(slide['content'], card_type)
                        slide['flashcards'] = parse_flashcards(raw, card_type)
                    except Exception as e:
                        st.error(f"Flashcard generation failed: {str(e)}")
                
                if slide['flashcards']:
                    cols = st.columns(2)
                    for i, card in enumerate(slide['flashcards']):
                        with cols[i % 2]:
                            with st.container():
                                st.markdown(f"""
                                <div style='padding:10px;border-radius:5px;
                                            background:#f0f2f6;margin:5px;'>
                                    <b>{card.get('question', card.get('term', 'Q'))}:</b><br>
                                    {card.get('options', '')}<br>
                                    <details style='margin-top:5px'>
                                        <summary>Reveal Answer</summary>
                                        <div style='padding:5px;background:white;
                                                    border-radius:3px;margin-top:5px'>
                                            {card.get('answer', card.get('definition', ''))}
                                        </div>
                                    </details>
                                </div>
                                """, unsafe_allow_html=True)
                    # Export buttons
                    df = pd.DataFrame(slide['flashcards'])
                    st.download_button(
                        label=f"Export Slide {slide['number']} Flashcards",
                        data=df.to_csv(index=False),
                        file_name=f"slide_{slide['number']}_cards.csv",
                        key=f"export_{idx}"
                    )
        
        progress_bar.progress((idx + 1) / len(slides))

# Chat feature
st.divider()
st.subheader("Presentation Chat Assistant 💬")

if "messages" not in st.session_state:
    st.session_state.messages = []

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

if prompt := st.chat_input("Ask about the presentation..."):
    context = " ".join([s['content'] for s in slides])
    try:
        response = safe_generate(f"Context: {context}\nQuestion: {prompt}\nAnswer:")
        st.session_state.messages.append({"role": "user", "content": prompt})
        st.session_state.messages.append({"role": "assistant", "content": response})
    except Exception as e:
        st.error(f"Chat error: {str(e)}")
    
    st.rerun()
